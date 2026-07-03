"""
AgentPath pitch deck generator.

Regenerate the deck after editing copy/colors below:
    python pitch/generate_deck.py

Output: pitch/AgentPath-Pitch.pptx  (16:9, 12 slides)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.line import MSO_LINE

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO = os.path.join(HERE, "assets", "logo.png")
OUT = os.path.join(HERE, "AgentPath-Pitch.pptx")

# ---------------------------------------------------------------- brand tokens
INK = RGBColor(0x37, 0x35, 0x2F)      # primary text (app ink)
MUTED = RGBColor(0x78, 0x77, 0x74)    # secondary text
FAINT = RGBColor(0x9B, 0x9A, 0x97)    # footnotes
BG = RGBColor(0xF6, 0xF5, 0xF4)       # app background
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE9, 0xE9, 0xE7)     # hairline borders
ACCENT = RGBColor(0x00, 0x75, 0xDE)   # logo blue
ACCENT_TINT = RGBColor(0xE9, 0xF2, 0xFC)
GREEN = RGBColor(0x2E, 0x7D, 0x32)

FONT = "Segoe UI"

SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.75)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

# ------------------------------------------------------------------- helpers


def slide_new(bg_color=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = bg_color
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def box(s, x, y, w, h, fill=CARD, line_color=LINE, radius=0.055):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, leading=1.0, space_after=0):
    """runs: str, or list of paragraphs; each paragraph is str or list of
    (text, {overrides}) tuples."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = leading
        p.space_after = Pt(space_after)
        if isinstance(para, str):
            para = [(para, {})]
        for txt, ov in para:
            r = p.add_run()
            r.text = txt
            f = r.font
            f.name = ov.get("font", FONT)
            f.size = Pt(ov.get("size", size))
            f.bold = ov.get("bold", bold)
            f.color.rgb = ov.get("color", color)
            if ov.get("italic"):
                f.italic = True
    return tb


def kicker(s, label, x=MARGIN, y=Inches(0.55)):
    text(s, x, y, Inches(6), Inches(0.35), label.upper(),
         size=13, color=ACCENT, bold=True)


def heading(s, title, x=MARGIN, y=Inches(0.95), w=None, size=33):
    text(s, x, y, w or (SW - 2 * MARGIN), Inches(1.1), title,
         size=size, bold=True, leading=1.05)


def footer(s, n):
    text(s, SW - Inches(1.2), SH - Inches(0.45), Inches(0.6), Inches(0.3),
         str(n), size=11, color=FAINT, align=PP_ALIGN.RIGHT)
    text(s, MARGIN, SH - Inches(0.45), Inches(3), Inches(0.3),
         "AgentPath", size=11, color=FAINT)


def image_zone(s, x, y, w, h, label):
    """Placeholder frame — drop a photo over it (or replace via script)."""
    z = box(s, x, y, w, h, fill=ACCENT_TINT, line_color=ACCENT, radius=0.035)
    z.line.dash_style = MSO_LINE.DASH
    z.line.width = Pt(1.25)
    text(s, x + Inches(0.3), y, w - Inches(0.6), h, [label],
         size=13, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def stat_card(s, x, y, w, h, big, label, source):
    box(s, x, y, w, h)
    pad = Inches(0.28)
    text(s, x + pad, y + Inches(0.22), w - 2 * pad, Inches(0.85),
         big, size=40, color=ACCENT, bold=True)
    text(s, x + pad, y + Inches(1.06), w - 2 * pad, h - Inches(1.4),
         label, size=13.5, color=INK, leading=1.12)
    text(s, x + pad, y + h - Inches(0.42), w - 2 * pad, Inches(0.3),
         source, size=10, color=FAINT)


def bullet_list(s, x, y, w, items, size=15, gap=0.5, dot_color=ACCENT):
    for i, (head, body) in enumerate(items):
        yy = y + Inches(i * gap)
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, x, yy + Inches(0.09),
                               Inches(0.09), Inches(0.09))
        d.fill.solid()
        d.fill.fore_color.rgb = dot_color
        d.line.fill.background()
        d.shadow.inherit = False
        runs = [[(head, {"bold": True})] + ([(" — " + body, {"color": MUTED})] if body else [])]
        text(s, x + Inches(0.25), yy, w - Inches(0.25), Inches(gap - 0.05),
             runs, size=size, leading=1.1)


# =========================================================== 1 · TITLE
s = slide_new(CARD)
# right half: image zone
image_zone(s, SW / 2 + Inches(0.2), Inches(0.6), SW / 2 - Inches(0.8),
           SH - Inches(1.2), "[ Image — a Nigerian student on campus ]")
# left: logo + name + one-liner
s.shapes.add_picture(LOGO, MARGIN, Inches(2.0), Inches(1.05), Inches(1.05))
text(s, MARGIN, Inches(3.15), Inches(5.6), Inches(1.0), "AgentPath",
     size=54, bold=True)
text(s, MARGIN, Inches(4.15), Inches(5.4), Inches(0.9),
     "An AI mentor for every first-generation student.",
     size=20, color=MUTED, leading=1.15)
chip = box(s, MARGIN, Inches(5.15), Inches(3.05), Inches(0.42),
           fill=ACCENT_TINT, line_color=None, radius=0.5)
text(s, MARGIN, Inches(5.15), Inches(3.05), Inches(0.42),
     "Agency-first  ·  Built for Nigeria", size=12.5, color=ACCENT, bold=True,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# =========================================================== 2 · PROBLEM
s = slide_new()
kicker(s, "The problem")
heading(s, "The students who need opportunities most\nare the least likely to hear about them.", size=29)
# left image zone
image_zone(s, MARGIN, Inches(2.35), Inches(4.5), Inches(4.35),
           "[ Image — student studying / campus life ]")
# right: three stat cards
cx = Inches(5.6)
cw = SW - cx - MARGIN
stat_card(s, cx, Inches(2.35), cw, Inches(1.32),
          "51%", "of the inequality in how far a child gets in school is explained by parental "
          "education — across nine Sub-Saharan countries incl. Nigeria.",
          "UNESCO Global Education Monitoring Report")
stat_card(s, cx, Inches(3.82), cw, Inches(1.32),
          "93%", "of Nigerian schools have no dedicated career counselor. First-gen students "
          "have no one at home to ask — and no one at school either.",
          "Teach For All")
stat_card(s, cx, Inches(5.29), cw, Inches(1.32),
          "1 : 1,000", "counselor-to-student ratio where a counselor exists at all. "
          "Guidance simply does not scale as people.",
          "African Journal of Career Development")
footer(s, 2)

# =========================================================== 3 · SOLUTION
s = slide_new()
kicker(s, "The solution")
heading(s, "A mentor in every pocket.")
# agency-first banner
bx, by, bw, bh = MARGIN, Inches(2.05), SW - 2 * MARGIN, Inches(1.0)
box(s, bx, by, bw, bh, fill=ACCENT, line_color=None)
text(s, bx + Inches(0.4), by, bw - Inches(0.8), bh,
     [[("Agency-first:  ", {"bold": True, "color": CARD}),
       ("the AI never opens with a list of options. It asks the student what they "
        "think and what they want — first, every time.", {"color": CARD})]],
     size=17, anchor=MSO_ANCHOR.MIDDLE, leading=1.15)
cards = [
    ("AI mentor chat", "A streaming conversational mentor that knows the student's "
     "course, year, and goals — and explains every acronym in plain language."),
    ("Curated opportunities", "A verified database of Nigerian scholarships, fellowships, "
     "internships, and competitions. The AI recommends only from it — no hallucinated deadlines."),
    ("Personal roadmap", "Every next step becomes a milestone the student owns and tracks — "
     "from “draft essay” to “submit application.”"),
]
cw = (SW - 2 * MARGIN - Inches(0.8)) / 3
for i, (t, d) in enumerate(cards):
    x = MARGIN + i * (cw + Inches(0.4))
    box(s, x, Inches(3.45), cw, Inches(3.0))
    text(s, x + Inches(0.3), Inches(3.75), cw - Inches(0.6), Inches(0.6),
         t, size=17, bold=True)
    text(s, x + Inches(0.3), Inches(4.4), cw - Inches(0.6), Inches(1.9),
         d, size=13.5, color=MUTED, leading=1.25)
footer(s, 3)

# =========================================================== 4 · PRODUCT
s = slide_new()
kicker(s, "Product")
heading(s, "Built and working today.")
# chat mock (left, dominant)
chx, chy, chw, chh = MARGIN, Inches(2.1), Inches(6.9), Inches(4.6)
box(s, chx, chy, chw, chh)
text(s, chx + Inches(0.35), chy + Inches(0.25), chw, Inches(0.35),
     "AI mentor — live chat", size=12, color=FAINT, bold=True)
# user bubble (right-aligned)
ub_w = Inches(4.4)
ub = box(s, chx + chw - ub_w - Inches(0.35), chy + Inches(0.75), ub_w, Inches(0.95),
         fill=ACCENT, line_color=None, radius=0.18)
text(s, chx + chw - ub_w - Inches(0.1), chy + Inches(0.75), ub_w - Inches(0.5), Inches(0.95),
     "I want to study abroad after graduation but I don't know where to start.",
     size=13.5, color=CARD, anchor=MSO_ANCHOR.MIDDLE, leading=1.15)
# ai bubble
ab_w = Inches(5.4)
box(s, chx + Inches(0.35), chy + Inches(1.95), ab_w, Inches(1.75),
    fill=BG, line_color=LINE, radius=0.12)
text(s, chx + Inches(0.6), chy + Inches(2.12), ab_w - Inches(0.5), Inches(1.5),
     "That's a solid goal, Amina. Before I point you anywhere — what matters "
     "most to you right now: funding, the country, or the course itself? "
     "What do you think?",
     size=13.5, color=INK, leading=1.2)
text(s, chx + Inches(0.35), chy + Inches(3.9), chw - Inches(0.7), Inches(0.5),
     [[("↑ The agency-first moment — the AI asks before it suggests.",
        {"italic": True, "color": ACCENT, "size": 12.5})]])
# right column: opportunities + roadmap mocks
rx = chx + chw + Inches(0.45)
rw = SW - rx - MARGIN
box(s, rx, Inches(2.1), rw, Inches(2.15))
text(s, rx + Inches(0.3), Inches(2.3), rw - Inches(0.6), Inches(0.35),
     "Opportunities", size=12, color=FAINT, bold=True)
for i, (t, tag) in enumerate([("MTN Foundation Scholarship", "Scholarship"),
                              ("NNPC/Seplat Undergraduate", "Scholarship"),
                              ("Google Africa Internship", "Internship")]):
    yy = Inches(2.7 + i * 0.48)
    text(s, rx + Inches(0.3), yy, rw - Inches(0.6), Inches(0.4),
         [[(t, {"size": 12, "bold": True})], ],)
    text(s, rx + Inches(0.3), yy + Inches(0.2), rw - Inches(0.6), Inches(0.25),
         tag, size=10, color=ACCENT)
box(s, rx, Inches(4.45), rw, Inches(2.25))
text(s, rx + Inches(0.3), Inches(4.65), rw - Inches(0.6), Inches(0.35),
     "Roadmap", size=12, color=FAINT, bold=True)
for i, (t, st, c) in enumerate([("Draft scholarship essay", "In progress", ACCENT),
                                ("Request transcript", "Pending", MUTED),
                                ("Submit MTN application", "Due Aug 15", GREEN)]):
    yy = Inches(5.05 + i * 0.5)
    text(s, rx + Inches(0.3), yy, rw - Inches(0.6), Inches(0.3),
         [[(t, {"size": 12, "bold": True})]])
    text(s, rx + Inches(0.3), yy + Inches(0.21), rw - Inches(0.6), Inches(0.25),
         st, size=10, color=c)
footer(s, 4)

# =========================================================== 5 · HOW IT WORKS
s = slide_new()
kicker(s, "How it works")
heading(s, "Four steps from lost to on-track.")
steps = [
    ("1 · Onboard", "Student shares university, course, year, and goals in a 3-step wizard."),
    ("2 · Talk", "The AI mentor engages — reflects goals back, asks what the student wants."),
    ("3 · Match", "Relevant opportunities surface from the curated database, in context."),
    ("4 · Act", "Next steps become roadmap milestones the student tracks to submission."),
]
cw = Inches(2.75)
gap = (SW - 2 * MARGIN - 4 * cw) / 3
for i, (t, d) in enumerate(steps):
    x = MARGIN + i * (cw + gap)
    box(s, x, Inches(2.5), cw, Inches(2.7))
    text(s, x + Inches(0.28), Inches(2.8), cw - Inches(0.56), Inches(0.5),
         t, size=17, bold=True, color=ACCENT)
    text(s, x + Inches(0.28), Inches(3.4), cw - Inches(0.56), Inches(1.7),
         d, size=13, color=MUTED, leading=1.25)
    if i < 3:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + cw + Inches(0.06),
                                Inches(3.68), gap - Inches(0.12), Inches(0.34))
        ar.fill.solid()
        ar.fill.fore_color.rgb = LINE
        ar.line.fill.background()
        ar.shadow.inherit = False
text(s, MARGIN, Inches(5.7), SW - 2 * MARGIN, Inches(0.5),
     [[("Under the hood:  ", {"bold": True})],
      [("Streaming AI grounded in a curated opportunity database — fast, cheap to run, "
        "and it never invents a deadline.", {"color": MUTED})]],
     size=14, leading=1.2)
footer(s, 5)

# =========================================================== 6 · WHY NOW
s = slide_new()
kicker(s, "Why now")
heading(s, "The market is huge. The moment is new.")
cw = (SW - 2 * MARGIN - Inches(0.8)) / 3
stats = [
    ("2.1M+", "students enrolled in Nigerian universities today — and growing every year.",
     "National Universities Commission"),
    ("~12%", "of university-age Nigerians reach higher education. Those who make it in "
     "carry their families' hopes — with the least guidance.",
     "UNESCO GEM Report"),
    ("1 in 2", "of 2M+ UTME candidates gained admission in 2025. Competition this fierce "
     "makes guidance decisive.", "JAMB, 2025"),
]
for i, (b, l, src) in enumerate(stats):
    x = MARGIN + i * (cw + Inches(0.4))
    stat_card(s, x, Inches(2.3), cw, Inches(2.6), b, l, src)
box(s, MARGIN, Inches(5.3), SW - 2 * MARGIN, Inches(1.0), fill=ACCENT_TINT, line_color=None)
text(s, MARGIN + Inches(0.4), Inches(5.3), SW - 2 * MARGIN - Inches(0.8), Inches(1.0),
     [[("Mentorship never scaled — people don't. ", {"bold": True}),
       ("LLM costs have collapsed: a full mentoring conversation now costs pennies. "
        "For the first time, one-on-one guidance can reach every student.",
        {"color": MUTED})]],
     size=15, anchor=MSO_ANCHOR.MIDDLE, leading=1.2)
footer(s, 6)

# =========================================================== 7 · COMPETITION
s = slide_new()
kicker(s, "Competition")
heading(s, "Everyone informs. Nobody mentors.")
rows = [
    ("", "Scholarship\naggregators", "Generic AI\nchatbots", "School\ncounselors", "AgentPath"),
    ("Curated, verified opportunities", "✓", "✗", "✗", "✓"),
    ("Conversational guidance", "✗", "✓", "✓", "✓"),
    ("Agency-first mentorship", "✗", "✗", "✓", "✓"),
    ("Scales to millions of students", "✓", "✓", "✗", "✓"),
]
tx, ty = MARGIN, Inches(2.35)
col_w = [Inches(4.3), Inches(1.9), Inches(1.9), Inches(1.9), Inches(1.9)]
row_h = Inches(0.82)
# highlight AgentPath column
hx = tx + sum(col_w[:4], Inches(0))
box(s, hx, ty - Inches(0.08), col_w[4], row_h * 5 + Inches(0.16),
    fill=ACCENT_TINT, line_color=ACCENT)
for r, row in enumerate(rows):
    y = ty + r * row_h
    for c, cell in enumerate(row):
        x = tx + sum(col_w[:c], Inches(0))
        is_head = r == 0
        is_mark = c > 0 and not is_head
        good = cell == "✓"
        text(s, x + Inches(0.1), y, col_w[c] - Inches(0.2), row_h,
             cell,
             size=13 if is_head else (17 if is_mark else 14),
             bold=is_head or (c == 4),
             color=(INK if is_head else (GREEN if good else FAINT)) if is_mark or is_head else INK,
             align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE, leading=1.05)
    if r < 4:
        ln = box(s, tx, ty + (r + 1) * row_h, sum(col_w, Inches(0)), Pt(1),
                 fill=LINE, line_color=None, radius=0)
footer(s, 7)

# =========================================================== 8 · BUSINESS MODEL
s = slide_new()
kicker(s, "Business model")
heading(s, "Free for students. Paid by those who want to reach them.")
cards = [
    ("Sponsored reach", "Foundations and companies already spend to find qualified "
     "applicants. AgentPath delivers matched, application-ready students.", "Revenue today"),
    ("Institutional partnerships", "Universities and NGOs license AgentPath as a "
     "guidance layer for their students and beneficiaries.", "Revenue next"),
    ("Premium mentorship", "Application review, deadline coaching, and interview prep "
     "for students who want to go deeper.", "Revenue later"),
]
cw = (SW - 2 * MARGIN - Inches(0.8)) / 3
for i, (t, d, tag) in enumerate(cards):
    x = MARGIN + i * (cw + Inches(0.4))
    box(s, x, Inches(2.4), cw, Inches(3.1))
    text(s, x + Inches(0.3), Inches(2.7), cw - Inches(0.6), Inches(0.35),
         tag.upper(), size=11, color=ACCENT, bold=True)
    text(s, x + Inches(0.3), Inches(3.1), cw - Inches(0.6), Inches(0.6),
         t, size=17, bold=True)
    text(s, x + Inches(0.3), Inches(3.75), cw - Inches(0.6), Inches(1.7),
         d, size=13, color=MUTED, leading=1.25)
text(s, MARGIN, Inches(5.85), SW - 2 * MARGIN, Inches(0.6),
     [[("The student never pays to discover an opportunity. ", {"bold": True}),
       ("That is non-negotiable — it is the reason the product exists.", {"color": MUTED})]],
     size=14.5)
footer(s, 8)

# =========================================================== 9 · TRACTION
s = slide_new()
kicker(s, "Where we are")
heading(s, "Built, not slideware.")
box(s, MARGIN, Inches(2.3), Inches(6.6), Inches(4.2))
text(s, MARGIN + Inches(0.35), Inches(2.6), Inches(6), Inches(0.4),
     "Working end-to-end today", size=15, bold=True)
bullet_list(s, MARGIN + Inches(0.35), Inches(3.2), Inches(6.0), [
    ("Full product live", "auth, onboarding, profiles"),
    ("Streaming AI mentor", "agency-first prompt, full conversation memory"),
    ("Curated opportunity database", "seeded with verified Nigerian programs"),
    ("Personal roadmap", "milestones linked to opportunities"),
    ("Production-grade stack", "NestJS + MongoDB + Next.js monorepo"),
], size=14.5, gap=0.62)
zx = MARGIN + Inches(7.0)
image_zone(s, zx, Inches(2.3), SW - zx - MARGIN, Inches(4.2),
           "[ Add early numbers here —\npilot users, waitlist signups,\npartner conversations ]")
footer(s, 9)

# =========================================================== 10 · GTM
s = slide_new()
kicker(s, "Go-to-market")
heading(s, "Campus-native distribution.")
items = [
    ("Student ambassadors", "Respected students on each campus onboard their departments — "
     "trust travels peer to peer."),
    ("WhatsApp-first sharing", "Departmental and faculty group chats are where opportunity "
     "links already spread. AgentPath links land where students already are."),
    ("Student unions & associations", "Official partnerships put AgentPath in orientation "
     "packs and union channels."),
    ("Scholarship foundations", "Providers want more qualified applicants — distribution "
     "is aligned, not adversarial. They promote AgentPath to their own audiences."),
]
for i, (t, d) in enumerate(items):
    x = MARGIN + (i % 2) * (Inches(6.0) + Inches(0.4))
    y = Inches(2.35) + (i // 2) * Inches(2.15)
    box(s, x, y, Inches(6.0), Inches(1.95))
    text(s, x + Inches(0.3), y + Inches(0.25), Inches(5.4), Inches(0.5),
         t, size=16, bold=True)
    text(s, x + Inches(0.3), y + Inches(0.8), Inches(5.4), Inches(1.0),
         d, size=13, color=MUTED, leading=1.25)
footer(s, 10)

# =========================================================== 11 · TEAM
s = slide_new()
kicker(s, "Team")
heading(s, "Why us.")
for i in range(2):
    x = MARGIN + i * Inches(6.4)
    box(s, x, Inches(2.4), Inches(5.6), Inches(3.2))
    av = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.35), Inches(2.75),
                            Inches(1.0), Inches(1.0))
    av.fill.solid()
    av.fill.fore_color.rgb = ACCENT_TINT
    av.line.color.rgb = ACCENT
    av.shadow.inherit = False
    text(s, x + Inches(1.55), Inches(2.85), Inches(3.8), Inches(0.45),
         "[ Founder name ]", size=17, bold=True)
    text(s, x + Inches(1.55), Inches(3.32), Inches(3.8), Inches(0.4),
         "[ Role — e.g. Founder & Engineer ]", size=13, color=ACCENT)
    text(s, x + Inches(0.35), Inches(4.05), Inches(4.9), Inches(1.4),
         "[ One line on why this person: built the product end-to-end / "
         "lived the first-gen journey / knows the campus ecosystem. ]",
         size=13, color=MUTED, leading=1.25)
text(s, MARGIN, Inches(5.95), SW - 2 * MARGIN, Inches(0.6),
     [[("If a founder is first-generation themselves — say it here, first. ",
        {"bold": True}),
       ("Lived experience is the strongest credential this deck can carry.",
        {"color": MUTED})]],
     size=14)
footer(s, 11)

# =========================================================== 12 · ASK
s = slide_new(INK)
text(s, MARGIN, Inches(1.0), Inches(8), Inches(0.4), "THE ASK",
     size=13, color=ACCENT, bold=True)
text(s, MARGIN, Inches(1.55), SW - 2 * MARGIN, Inches(1.1),
     "[ Your ask — e.g. $150k pre-seed / program acceptance ]",
     size=34, bold=True, color=CARD)
asks = [
    ("Pilot on 3 campuses", "ambassador programs at three universities"),
    ("Reach 10,000 students", "onboarded and actively mentored in year one"),
    ("Grow the database", "verified opportunities, refreshed weekly"),
]
cw = (SW - 2 * MARGIN - Inches(0.8)) / 3
for i, (t, d) in enumerate(asks):
    x = MARGIN + i * (cw + Inches(0.4))
    box(s, x, Inches(3.1), cw, Inches(1.8), fill=RGBColor(0x45, 0x43, 0x3C),
        line_color=None)
    text(s, x + Inches(0.3), Inches(3.4), cw - Inches(0.6), Inches(0.5),
         t, size=16, bold=True, color=CARD)
    text(s, x + Inches(0.3), Inches(3.95), cw - Inches(0.6), Inches(0.8),
         d, size=12.5, color=RGBColor(0xBD, 0xBC, 0xB8), leading=1.2)
text(s, MARGIN, Inches(5.6), SW - 2 * MARGIN, Inches(1.2),
     [[("Every student deserves someone who asks: ", {"color": CARD}),
       ("“what do you think?”", {"color": ACCENT, "bold": True})]],
     size=24, leading=1.15)
s.shapes.add_picture(LOGO, SW - Inches(1.45), SH - Inches(1.45),
                     Inches(0.7), Inches(0.7))

prs.save(OUT)
print(f"Saved {OUT} — {len(prs.slides.slides if hasattr(prs.slides, 'slides') else prs.slides._sldIdLst)} slides")
